#!/usr/bin/env python3
"""
qa-validate.py <file.pptx>

Mandatory final step before shipping the presentation (see DESIGN_SYSTEM.md /
Preschool_Rulebook.md, QA validator section):
1. schema validation (its own checks via python-pptx, needs nothing external)
2. if LibreOffice is installed - render to PDF + pdfplumber: the real
   bounding boxes of words and images must not go past the page edges
   (textbox overflow isn't caught by schema validation, only by rendering).
   If LibreOffice isn't found, this step is skipped with an explanation (not an error).
3. a simple check for em dashes and common typo/placeholder text
   (via the markitdown library, no external command-line calls).

Exits non-zero if problems are found.
"""
import sys
import subprocess
import shutil
import os
import re
import json

# guard against crashing on Windows consoles with non-UTF8 encodings (cp1252 etc.)
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def lint_content(content_path, problems, warnings):
    """Checks on content.json itself, besides the visual layout."""
    with open(content_path, encoding="utf-8") as f:
        content = json.load(f)

    qa = content.get("presentation_qa", [])
    if len(qa) < 4:
        problems.append(f"content.json: presentation_qa has only {len(qa)} question(s), minimum is 4.")

    props = content.get("story_props", [])
    if len(props) > 4:
        problems.append(f"content.json: story_props has {len(props)} items, keep it minimal (max ~4).")
    for prop in props:
        if prop.get("role") not in ("character", "handout"):
            problems.append(
                f"content.json: story_props item '{prop.get('name', '?')}' is missing a valid 'role' "
                f"(character/handout) - printing can't size it correctly without this."
            )
    game2 = content.get("game2", {})
    if game2.get("type") == "search" and game2.get("search_item") and not game2["search_item"].get("image_prompt"):
        problems.append("content.json: game2.search_item is set but missing image_prompt.")

    # Print-asset completeness: catches the "duck sorting game with no duck
    # image" bug at generation time instead of at print time. print_items is
    # the generic mechanism for matching/sorting/compare/pattern; colors is
    # the older color-recolor mechanism; matching/search may legitimately
    # reuse already-printed story_props instead (the guide's documented
    # fallback), so that combination only warns, never fails.
    print_items = game2.get("print_items", [])
    has_colors = bool(game2.get("colors"))
    has_print_assets = bool(print_items) or has_colors

    for i, item in enumerate(print_items):
        name = item.get("name")
        if not name:
            problems.append(f"content.json: game2.print_items[{i}] is missing 'name'.")
        if not item.get("image_prompt"):
            problems.append(f"content.json: game2.print_items[{i}] ('{name or '?'}') is missing 'image_prompt'.")
        copies = item.get("copies")
        if not isinstance(copies, int) or isinstance(copies, bool) or copies < 1:
            problems.append(f"content.json: game2.print_items[{i}] ('{name or '?'}') needs a positive integer 'copies', got {copies!r}.")

    game2_type = game2.get("type")
    if game2_type in ("sorting", "pattern") and not has_print_assets:
        problems.append(
            f"content.json: game2.type is '{game2_type}' but neither print_items nor colors is set - "
            f"nothing will be printed for this game (this is the exact bug that shipped a duck-sorting "
            f"game with no duck images). Add game2.print_items."
        )
    if game2_type == "compare" and not has_print_assets:
        problems.append(
            "content.json: game2.type is 'compare' but no print_items/colors is set - 'compare' has no "
            "story-prop-reuse fallback, it always needs its own print_items (typically 2 entries, copies: 1 each)."
        )
    if game2_type == "matching" and not has_print_assets:
        warnings.append(
            "content.json: game2.type is 'matching' with no print_items/colors - only OK if this "
            "intentionally reuses the already-printed story_props cards (per the guide); if not, add print_items."
        )

    if game2_type in ("sorting", "matching") and print_items:
        total = sum(int(it.get("copies") or 0) for it in print_items)
        if total < 4 or total > 12:
            warnings.append(
                f"content.json: game2.print_items total copies is {total} for a '{game2_type}' game - "
                f"usually you want about 8 total (one class set), see CONTENT_GENERATION_GUIDE.md Game 2 print quantity rule."
            )
    if game2_type == "compare":
        for item in print_items:
            copies = item.get("copies")
            if isinstance(copies, int) and copies > 2:
                warnings.append(
                    f"content.json: game2.print_items '{item.get('name')}' has copies={copies} for a 'compare' game - "
                    f"compare items are held up by the teacher, not handed out; 1-2 copies is usually enough."
                )

    for prop in props:
        if "class_copies" in prop:
            class_copies = prop["class_copies"]
            if prop.get("role") != "handout":
                problems.append(f"content.json: story_props item '{prop.get('name', '?')}' has 'class_copies' but role is '{prop.get('role')}' - class_copies only applies to handout props.")
            elif not isinstance(class_copies, int) or isinstance(class_copies, bool) or class_copies < 1:
                problems.append(f"content.json: story_props item '{prop.get('name', '?')}' has an invalid class_copies ({class_copies!r}), must be a positive integer.")
            elif class_copies < 4:
                warnings.append(f"content.json: story_props item '{prop.get('name')}' has class_copies={class_copies} - too few for a full class if each child gets one; consider 8.")

    for game_key in ("game1", "game2", "game3"):
        game = content.get(game_key, {})
        script = game.get("script")
        if not script:
            problems.append(f"content.json: {game_key}.script is missing or empty (unified script format required).")

    topic = content.get("topic", "").lower()
    summary = content.get("story", {}).get("short_summary", "").lower()
    if topic and topic not in summary:
        problems.append(
            f"content.json: WARNING — story.short_summary does not mention the topic word "
            f"'{content.get('topic')}'. Double-check the story leads to building the {content.get('topic')} "
            f"model itself, not a derivative object (a derivative like a home/container belongs in Challenge)."
        )


def check_video_availability(content_path, problems, warnings):
    """
    Checks that every youtube id used in the lesson is actually reachable
    via the public oEmbed endpoint - no API key needed. Requires internet.
    404/401 -> hard error, anything else -> warning (could be a network
    restriction/rate limiting rather than a real problem with the video).
    """
    try:
        import urllib.request
        import urllib.error
    except ImportError:
        warnings.append("Video check skipped: urllib not available.")
        return

    with open(content_path, encoding="utf-8") as f:
        content = json.load(f)

    ids_to_check = {}
    if content.get("intro_video_youtube_id"):
        ids_to_check["intro_video_youtube_id"] = content["intro_video_youtube_id"]

    ids_to_check["alphabet (constant)"] = "-1jxqVy5SlA"
    ids_to_check["closing (constant)"] = "h2AndZKYZBQ"

    letter = content.get("letter", "").upper()
    letter_videos_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "letter-videos.json")
    if os.path.exists(letter_videos_path):
        with open(letter_videos_path, encoding="utf-8") as f:
            letter_videos = json.load(f)
        vid = letter_videos.get(letter)
        if vid:
            ids_to_check[f"letter-specific ({letter})"] = vid

    for label, vid in ids_to_check.items():
        url = f"https://www.youtube.com/oembed?url=https://www.youtube.com/watch?v={vid}&format=json"
        try:
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=8) as resp:
                if resp.status != 200:
                    warnings.append(f"Video check: {label} (id={vid}) returned status {resp.status} — may be unavailable.")
        except urllib.error.HTTPError as e:
            if e.code in (404, 401):
                problems.append(f"Video UNAVAILABLE: {label} (id={vid}) — oEmbed returned {e.code} (video deleted/private/nonexistent). Replace this youtube id.")
            else:
                warnings.append(f"Video check inconclusive for {label} (id={vid}): oEmbed returned {e.code}. Often network restriction/rate limiting rather than a bad video — verify manually at https://youtube.com/watch?v={vid} if unsure.")
        except Exception as e:
            warnings.append(f"Video check: {label} (id={vid}) could not be verified ({e}). Check manually or verify internet access.")


def schema_validate(pptx_path, problems):
    """Self-contained slide-bounds check via python-pptx - needs no
    external scripts/paths, works on any machine. Only text and images are
    checked strictly - decorative background shapes with no text (e.g. the
    corner squares) may intentionally poke slightly past the slide edge,
    that's a deliberate design choice, not a bug."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    tolerance = 9525 * 2  # ~2pt margin for rounding

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
            except (AttributeError, TypeError):
                continue
            if left is None or top is None or width is None or height is None:
                continue

            has_text = bool(getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            if not has_text and not is_picture:
                continue  # purely decorative shape - bleeding past the edge is fine

            if left < -tolerance or top < -tolerance:
                problems.append(f"Slide {i}: shape starts outside left/top bounds.")
            if left + width > sw + tolerance:
                problems.append(f"Slide {i}: shape overflows right edge (x+w={left+width}, slide width={sw}).")
            if top + height > sh + tolerance:
                problems.append(f"Slide {i}: shape overflows bottom edge (y+h={top+height}, slide height={sh}).")


def find_soffice():
    """Looks for LibreOffice on the machine under various names/paths. None if not found."""
    for name in ("soffice", "libreoffice", "soffice.exe"):
        path = shutil.which(name)
        if path:
            return path
    # common install paths on Windows
    common_win_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in common_win_paths:
        if os.path.exists(p):
            return p
    return None


def bounding_box_check(pptx_path, problems):
    """If LibreOffice is installed, renders to PDF and checks the real
    bounding boxes of text/images. If LibreOffice isn't found, skips this
    step with an explanation - this is NOT an error (just a less thorough check)."""
    soffice = find_soffice()
    if not soffice:
        print("LibreOffice not found on this computer - skipping the text/image overflow check.")
        print("(This is not an error. Open the presentation in PowerPoint and check the")
        print(" slides by eye, especially whether text overlaps neighboring elements.)")
        return

    base = os.path.splitext(pptx_path)[0]
    pdf_path = base + ".pdf"
    r = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", pptx_path],
        capture_output=True, text=True, cwd=os.path.dirname(pptx_path),
    )
    print(r.stdout, r.stderr)

    try:
        import pdfplumber
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "pdfplumber", "-q"])
        import pdfplumber

    if not os.path.exists(pdf_path):
        problems.append(f"PDF render not found at {pdf_path}, could not run bounding-box QA.")
        return

    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            pw, ph = page.width, page.height
            for w in page.extract_words():
                if w["x0"] < -1 or w["x1"] > pw + 1 or w["top"] < -1 or w["bottom"] > ph + 1:
                    problems.append(
                        f"Slide {i}: text '{w['text']}' overflows page bounds "
                        f"(x0={w['x0']:.1f}, x1={w['x1']:.1f}, page width={pw:.1f})"
                    )
            for im in page.images:
                if im["x0"] < -1 or im["x1"] > pw + 1 or im["top"] < -1 or im["bottom"] > ph + 1:
                    problems.append(f"Slide {i}: image overflows page bounds")


def text_hygiene_check(pptx_path, problems):
    """Uses the markitdown library directly (no external command-line call -
    this is more reliable on Windows, where a script's .exe might not be
    on PATH)."""
    try:
        from markitdown import MarkItDown
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "markitdown[pptx]", "-q"])
        from markitdown import MarkItDown

    try:
        md = MarkItDown()
        result = md.convert(pptx_path)
        text = result.text_content
    except Exception as e:
        print(f"Text hygiene check skipped (could not read text: {e})")
        return

    if "\u2014" in text:
        problems.append("Found an em dash (\u2014) — replace with comma/colon per style rules.")
    if re.search(r"\bTODO\b|\[insert|xxx", text, re.IGNORECASE):
        problems.append("Found placeholder text (TODO / [insert / xxx) left in the deck.")


def check_image_cropping(content_path, problems):
    """Heuristic crop detector: for a properly generated cutout (transparent
    background), the real subject should never touch the image edge - if it
    does, the photo/illustration was cropped (e.g. ears cut off). Checks the
    alpha channel at the four borders; a large opaque fraction on any edge
    means the subject likely runs off the frame."""
    try:
        from PIL import Image
    except ImportError:
        print("SKIP: Pillow not installed, cannot check for cropped images")
        return

    with open(content_path, "r", encoding="utf-8") as f:
        content = json.load(f)
    slug = content["topic"].lower().replace(" ", "_")
    data_dir = os.environ.get("BRICK_DATA_DIR", os.path.join(os.path.dirname(__file__), ".."))
    gen_dir = os.path.join(data_dir, "assets", "generated", slug)

    targets = ["real_object.png"]
    for prop in content.get("story_props", []):
        targets.append(prop["name"].lower().replace(" ", "_") + ".png")

    THRESHOLD = 0.12  # 12%+ opaque pixels touching an edge -> likely cropped

    for fname in targets:
        img_path = os.path.join(gen_dir, fname)
        if not os.path.exists(img_path):
            continue
        try:
            img = Image.open(img_path).convert("RGBA")
        except Exception:
            continue
        w, h = img.size
        px = img.load()

        def edge_ratio(coords):
            if not coords:
                return 0
            opaque = sum(1 for (x, y) in coords if px[x, y][3] > 200)
            return opaque / len(coords)

        step = max(1, w // 200)
        top = [(x, 0) for x in range(0, w, step)]
        bottom = [(x, h - 1) for x in range(0, w, step)]
        step_h = max(1, h // 200)
        left = [(0, y) for y in range(0, h, step_h)]
        right = [(w - 1, y) for y in range(0, h, step_h)]

        worst = max(edge_ratio(top), edge_ratio(bottom), edge_ratio(left), edge_ratio(right))
        if worst > THRESHOLD:
            problems.append(
                f"{fname}: subject appears to touch the image edge ({worst*100:.0f}% of an edge is opaque) "
                f"- likely cropped (e.g. ears/tail cut off). Regenerate with a prompt that leaves clear margin "
                f"around the full subject."
            )


def main():
    if len(sys.argv) < 2:
        print("Usage: qa-validate.py <file.pptx> [content.json]")
        sys.exit(1)

    pptx_path = os.path.abspath(sys.argv[1])
    problems = []
    warnings = []

    if len(sys.argv) >= 3:
        content_path = os.path.abspath(sys.argv[2])
        print("== 0/4 Content lint ==")
        lint_warnings = []
        lint_content(content_path, problems, lint_warnings)
        for p in problems:
            print("LINT:", p)
        for w in lint_warnings:
            print("LINT WARNING:", w)
        print("== 0b/4 Video availability (youtube oEmbed, needs internet) ==")
        check_video_availability(content_path, problems, warnings)
        for p in warnings:
            print("VIDEO WARNING:", p)
        for p in [x for x in problems if x.startswith("Video UNAVAILABLE")]:
            print("VIDEO FAIL:", p)
        print("== 0c/4 Image crop check ==")
        check_image_cropping(content_path, problems)

    print("== 1/3 Schema validation (bounds check) ==")
    schema_validate(pptx_path, problems)

    print("== 2/3 Rendering to PDF for bounding-box check ==")
    bounding_box_check(pptx_path, problems)

    print("== 3/3 Text hygiene (em-dash, placeholders) ==")
    text_hygiene_check(pptx_path, problems)

    print("\n==================== QA SUMMARY ====================")
    if problems:
        for p in problems:
            print("FAIL:", p)
        print(f"\n{len(problems)} problem(s) found. Fix before sending to the user.")
        sys.exit(1)
    else:
        print("All QA checks passed.")
        sys.exit(0)


if __name__ == "__main__":
    main()
